"""
PPT解析Agent - 用于解析课件PPT
支持文本提取、表格提取、图表识别、布局分析
"""
import logging
from typing import Dict, Any, List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTParserAgent:
    """PowerPoint文档解析Agent"""

    def __init__(self):
        """初始化PPT解析器"""
        self.name = "PPT Parser"
        self.description = "解析PowerPoint课件，提取文本、表格、图表和结构信息"
        logger.info(f"{self.name} initialized")

    def parse_ppt(self, ppt_path: str, slide_range: Optional[tuple] = None) -> Dict[str, Any]:
        """
        解析PPT文件

        Args:
            ppt_path: PPT文件路径
            slide_range: 幻灯片范围 (start, end)，None表示全部

        Returns:
            包含解析结果的字典
        """
        try:
            prs = Presentation(ppt_path)
            total_slides = len(prs.slides)

            # 确定要解析的幻灯片范围
            if slide_range:
                start_slide, end_slide = slide_range
                start_slide = max(0, start_slide)
                end_slide = min(total_slides, end_slide)
            else:
                start_slide, end_slide = 0, total_slides

            logger.info(f"Parsing PPT: {ppt_path}, slides {start_slide}-{end_slide} of {total_slides}")

            # 解析结果
            result = {
                "file_name": Path(ppt_path).name,
                "total_slides": total_slides,
                "parsed_slides": end_slide - start_slide,
                "slides": [],
                "metadata": self._extract_metadata(prs),
                "structure": {
                    "sections": [],
                    "main_topics": []
                }
            }

            # 逐页解析
            for slide_num in range(start_slide, end_slide):
                slide_data = self._parse_slide(prs.slides[slide_num], slide_num + 1)
                result["slides"].append(slide_data)

                # 提取章节标题
                if slide_data.get("title"):
                    result["structure"]["sections"].append({
                        "slide": slide_num + 1,
                        "title": slide_data["title"]
                    })

            # 分析文档结构
            result["structure"]["main_topics"] = self._extract_main_topics(result["slides"])

            logger.info(f"PPT parsing completed: {len(result['slides'])} slides processed")
            return result

        except Exception as e:
            logger.error(f"Error parsing PPT: {str(e)}")
            raise

    def _parse_slide(self, slide, slide_num: int) -> Dict[str, Any]:
        """
        解析单个幻灯片

        Args:
            slide: python-pptx幻灯片对象
            slide_num: 幻灯片编号

        Returns:
            幻灯片解析结果
        """
        # 提取标题
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()

        # 提取文本内容
        text_frames = []
        tables = []
        charts = []
        images = []

        for shape in slide.shapes:
            # 文本框
            if shape.has_text_frame:
                text_content = self._extract_text_from_shape(shape)
                if text_content and shape != slide.shapes.title:
                    text_frames.append(text_content)

            # 表格
            elif shape.has_table:
                table_data = self._extract_table(shape.table)
                tables.append(table_data)

            # 图表
            elif shape.has_chart:
                chart_info = self._extract_chart_info(shape)
                charts.append(chart_info)

            # 图片
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_info = self._extract_image_info(shape, slide_num)
                images.append(image_info)

        # 合并所有文本
        all_text = title + "\n" + "\n".join(text_frames)

        # 提取公式和关键词
        formulas = self._extract_formulas(all_text)
        keywords = self._extract_keywords(all_text)

        return {
            "slide_num": slide_num,
            "title": title,
            "text_frames": text_frames,
            "full_text": all_text.strip(),
            "tables": tables,
            "charts": charts,
            "images": images,
            "formulas": formulas,
            "keywords": keywords,
            "layout": self._detect_layout_type(slide)
        }

    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """提取PPT元数据"""
        core_props = prs.core_properties

        return {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height
        }

    def _extract_text_from_shape(self, shape) -> str:
        """从形状中提取文本"""
        text_parts = []
        for paragraph in shape.text_frame.paragraphs:
            para_text = paragraph.text.strip()
            if para_text:
                text_parts.append(para_text)
        return "\n".join(text_parts)

    def _extract_table(self, table) -> Dict[str, Any]:
        """提取表格数据"""
        rows = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            rows.append(row_data)

        return {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "data": rows,
            "header": rows[0] if rows else []
        }

    def _extract_chart_info(self, shape) -> Dict[str, Any]:
        """提取图表信息"""
        chart = shape.chart

        return {
            "type": str(chart.chart_type),
            "title": chart.chart_title.text_frame.text if chart.has_title else "",
            "series_count": len(chart.series),
            # 可以进一步提取数据点，但这里保持简单
        }

    def _extract_image_info(self, shape, slide_num: int) -> Dict[str, Any]:
        """提取图片信息"""
        return {
            "slide": slide_num,
            "name": shape.name,
            "width": shape.width,
            "height": shape.height,
            "left": shape.left,
            "top": shape.top,
        }

    def _detect_layout_type(self, slide) -> str:
        """检测幻灯片布局类型"""
        # 根据形状数量和类型推断布局
        shape_count = len(slide.shapes)

        has_title = slide.shapes.title is not None
        has_table = any(s.has_table for s in slide.shapes)
        has_chart = any(s.has_chart for s in slide.shapes)
        has_image = any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes)

        if has_title and shape_count <= 2:
            return "title_only"
        elif has_table:
            return "table"
        elif has_chart:
            return "chart"
        elif has_image:
            return "image"
        elif has_title and shape_count > 2:
            return "title_content"
        else:
            return "blank"

    def _extract_formulas(self, text: str) -> List[str]:
        """识别文本中的数学公式"""
        formulas = []

        # 数学公式的常见模式
        formula_patterns = [
            r'[a-zA-Z]\s*[=+\-*/]\s*[a-zA-Z0-9()]+',
            r'\b[a-zA-Z]\s*\^\s*[0-9]+',
            r'∫|∑|∏|∂|∇|√|π|θ|α|β|γ|λ|μ|σ|Δ',
            r'\b(sin|cos|tan|log|ln|exp|lim|max|min|det|rank)\s*\(',
            r'[a-zA-Z]+\s*=\s*\{[^}]+\}',
            r'\([^)]*[+\-*/^][^)]*\)',
        ]

        lines = text.split('\n')
        for line in lines:
            for pattern in formula_patterns:
                matches = re.findall(pattern, line)
                if matches:
                    cleaned_line = line.strip()
                    if cleaned_line and len(cleaned_line) < 200:
                        formulas.append(cleaned_line)
                    break

        return list(set(formulas))[:15]

    def _extract_keywords(self, text: str) -> List[str]:
        """提取关键词"""
        # 专业领域关键词库
        keywords_db = [
            # 数学
            '定理', '公式', '证明', '推导', '定义', '性质', '引理', '推论',
            '矩阵', '向量', '函数', '导数', '积分', '微分', '极限', '级数',
            '概率', '统计', '方程', '不等式', '集合', '映射', '变换',
            # 经济学
            '需求', '供给', '均衡', '弹性', '效用', '成本', '收益', '利润',
            '市场', '价格', '产量', '消费', '投资', '储蓄', 'GDP',
            # 计算机
            '算法', '复杂度', '数据结构', '时间复杂度', '空间复杂度',
            '排序', '搜索', '树', '图', '哈希', '递归', '动态规划',
            # 英文
            'theorem', 'proof', 'definition', 'matrix', 'vector',
            'derivative', 'integral', 'algorithm', 'complexity'
        ]

        found_keywords = []
        text_lower = text.lower()

        for keyword in keywords_db:
            if keyword.lower() in text_lower:
                found_keywords.append(keyword)

        return list(set(found_keywords))[:12]

    def _extract_main_topics(self, slides: List[Dict]) -> List[str]:
        """从所有幻灯片中提取主要主题"""
        all_keywords = []
        for slide in slides:
            all_keywords.extend(slide.get("keywords", []))

        # 统计关键词频率
        from collections import Counter
        keyword_freq = Counter(all_keywords)

        # 返回出现频率最高的主题
        main_topics = [kw for kw, _ in keyword_freq.most_common(8)]
        return main_topics

    def extract_section(self, ppt_path: str, section_title: str) -> Dict[str, Any]:
        """
        提取特定章节的幻灯片

        Args:
            ppt_path: PPT文件路径
            section_title: 章节标题（用于匹配）

        Returns:
            章节内容
        """
        # 首先解析整个文档
        doc_data = self.parse_ppt(ppt_path)

        # 查找匹配的章节
        section_slides = []
        in_section = False

        for i, slide in enumerate(doc_data["slides"]):
            # 检查是否是目标章节的开始
            if section_title.lower() in slide.get("title", "").lower():
                in_section = True
                section_slides.append(slide)
            # 如果遇到新章节标题，停止
            elif in_section and slide.get("title") and slide.get("layout") == "title_only":
                break
            # 在章节内，继续收集
            elif in_section:
                section_slides.append(slide)

        return {
            "section_title": section_title,
            "slides": section_slides,
            "total_slides": len(section_slides),
            "text": "\n\n".join(slide["full_text"] for slide in section_slides),
            "formulas": [f for slide in section_slides for f in slide.get("formulas", [])],
            "keywords": list(set([k for slide in section_slides for k in slide.get("keywords", [])]))
        }
